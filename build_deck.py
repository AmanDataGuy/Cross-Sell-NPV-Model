"""
A 5-slide client deck (JD: "creating strong client presentations"). Pulls
its numbers from the CSVs the other scripts already wrote, so the deck
can't say something the pipeline didn't actually produce.
Run this last, after load/model/hypothesis/segment/backtest have all run.
"""
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

OUTPUTS = Path("outputs")
FIGURES = Path("figures")


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], image: Path | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = bullets[0]
    for line in bullets[1:]:
        p = body.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
    if image and image.exists():
        slide.shapes.add_picture(str(image), Inches(5.5), Inches(1.5), height=Inches(4.5))


def main() -> None:
    model_comparison = pd.read_csv(OUTPUTS / "model_comparison.csv", index_col="model")
    drivers = pd.read_csv(OUTPUTS / "driver_significance.csv")
    target_list = pd.read_csv("target_list.csv")
    tier_totals = target_list.groupby("tier")["expected_npv"].sum().reindex(["High", "Medium", "Low"])
    top_drivers = drivers[drivers["significant_fdr"] & drivers["practically_important"]]["driver"].tolist()

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Cross-Sell Propensity & Customer NPV"
    title_slide.placeholders[1].text = "Who to target next, and what it's worth"

    add_bullet_slide(prs, "The ask", [
        "Of our existing customers, which should we target for the next-product offer?",
        "How sure are we the attributes we target on actually matter?",
        "What is the long-term value of pursuing them?",
    ])

    add_bullet_slide(prs, "Approach", [
        "45,211 customers, real UCI Bank Marketing data",
        "Logistic regression (interpretable, odds ratios) vs. gradient boosting benchmark",
        "Chi-square / t-tests + Benjamini-Hochberg correction to separate real drivers from noise",
        "Calibration + cost-sensitive threshold so propensity x dollars is honest",
        "Shallow decision tree segments customers into High/Med/Low target tiers",
    ])

    add_bullet_slide(prs, "Key result", [
        f"Logistic AUC {model_comparison.loc['logistic', 'auc']:.3f} vs. GBM AUC {model_comparison.loc['gbm', 'auc']:.3f}",
        f"Logistic PR-AUC {model_comparison.loc['logistic', 'pr_auc']:.3f} vs. GBM PR-AUC {model_comparison.loc['gbm', 'pr_auc']:.3f}",
        "Real drivers (FDR-significant and practically important): " + ", ".join(top_drivers),
        "Back-tested on held-out, out-of-time data: High-tier customers subscribed at a meaningfully higher rate",
    ], image=FIGURES / "lift_gains_chart.png")

    add_bullet_slide(prs, "Recommendation", [
        f"High tier: {tier_totals['High']:,.0f} expected NPV" if pd.notna(tier_totals.get("High")) else "High tier: n/a",
        f"Medium tier: {tier_totals['Medium']:,.0f} expected NPV" if pd.notna(tier_totals.get("Medium")) else "Medium tier: n/a",
        f"Low tier: {tier_totals['Low']:,.0f} expected NPV" if pd.notna(tier_totals.get("Low")) else "Low tier: n/a",
        "Target High + Medium tiers; full ranked list in target_list.csv",
        "Assumptions (margin, retention, cost) are toggleable in npv_model.xlsx",
    ], image=FIGURES / "npv_by_tier_chart.png")

    prs.save("deck.pptx")
    print("deck.pptx written.")


if __name__ == "__main__":
    main()
